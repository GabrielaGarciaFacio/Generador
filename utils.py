import os
import sys
import requests
import io
import random
from io import BytesIO
import pyodbc
import requests
import tempfile
import openpyxl
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from copy import copy
import pdfplumber
import math
from pptx import Presentation
from docx import Document
import time
import traceback
import streamlit as st
import getpass
import requests
import re

def extract_text_from_pdf_bdd(clave):
    """
    Extrae texto de un archivo PDF.
    
    :param clave: Clave que proporciona el usuario para obtener el PDF
    :return: Todo el texto extraído como una cadena.
    """
    connectionString = ("DRIVER={ODBC Driver 18 for SQL Server};"
                            "SERVER=scenetecprod.czbotsckvb07.us-west-2.rds.amazonaws.com;"
                            "DATABASE=netec_prod;"
                            "UID=netec_read;"
                            "PWD=R3ad25**SC3.2025-;"
                            "TrustServerCertificate=yes;")
    text = ""
    try:
        #Conectar a la base de datos            
        conn=pyodbc.connect(connectionString)
        cursor=conn.cursor()

        #Consultar el PDF almacenado en la BD usando la clave proporcionada
        query="""
            SELECT link_temario 
            FROM cursos_habilitados 
            WHERE link_temario IS NOT NULL
                AND link_temario <> ''
                AND link_temario <> 'NA'
                AND clave=?
            """
        
        print(f"[DEBUG] Clave usada en query: '{clave}'")  # para detectar espacios invisibles

        cursor.execute(query,(clave,))

        #Obtener el PDF en formato binario
        pdf_data=cursor.fetchone()
        if pdf_data:
            #Base URL
            base_url="https://sce.netec.com/"
            #Construir el enlace completo
            full_link=base_url+pdf_data[0]
    
            try:
                response=requests.get(full_link)
                response.raise_for_status() #verifica si la descarga fue exitosa
                pdf_bytes=io.BytesIO(response.content)

                #Extraer texto del PDF
                text=""
                with pdfplumber.open(pdf_bytes) as pdf:
                    for page in pdf.pages:
                        text += page.extract_text() + "\n"  # Agrega un salto de línea entre páginas
                text=f"source: {clave}\n"+ text
                return text
            except requests.exceptions.RequestException as e:
                print(f"Error al descargar el PDF: {e}")
                return ""
            except Exception as e:
                print(f'Erros al extraer texto del PDF: {e}')
                return ""
        else:
            print(f'No se encontró el PDF para la clave {clave}')
            return ""

    except Exception as e:
        print(f'Error al conectar la base de datos o ejecutar la consulta: {e}')
        return ""

def apply_format(worksheet, source_row, target_row):
    """
    Aplica el formato de una fila fuente (source_row) a una fila objetivo (target_row).
    
    :param worksheet: Worksheet del archivo Excel.
    :param source_row: Número de la fila con el formato base.
    :param target_row: Número de la fila a la que se aplicará el formato.
    """
    for col in range(1, worksheet.max_column + 1):  # Iterar por todas las columnas
        source_cell = worksheet.cell(row=source_row, column=col)
        target_cell = worksheet.cell(row=target_row, column=col)
        
        # Copiar valor y estilo
        target_cell.font = copy(source_cell.font)
        target_cell.border = copy(source_cell.border)
        target_cell.fill = copy(source_cell.fill)
        target_cell.number_format = copy(source_cell.number_format)
        target_cell.protection = copy(source_cell.protection)
        target_cell.alignment = copy(source_cell.alignment)

def resource_path(relative_path):
    # Obtiene la ruta absoluta al recurso, tanto si estamos en desarrollo como empaquetado
    if hasattr(sys, '_MEIPASS'):
        return os.path.join(sys._MEIPASS, relative_path)
    return os.path.abspath(relative_path)
def translate_with_mymemory(text):
    """
    Traduce texto del español al inglés utilizando el servicio gratuito MyMemory Translation API.
    El texto se divide en fragmentos significativos para mantener la coherencia.
    
    :param text: Texto a traducir (string)
    :return: Texto traducido (string)
    """
    import requests
    import time
    import re
    
    # Si el texto está vacío, devolver vacío
    if not text or text.strip() == "":
        return ""
    
    # Dividir el texto en secciones significativas
    # Cada sección comienza con "Tema:" o "Pregunta:"
    sections = []
    current_section = ""
    
    # Dividir por líneas
    lines = text.split('\n')
    
    for line in lines:
        if line.strip().startswith("TEMA:") or re.match(r'^\d+\.\s*PREGUNTA:', line):
            # Si ya hay una sección acumulada, la añadimos a la lista
            if current_section.strip():
                sections.append(current_section.strip())
            # Iniciamos una nueva sección
            current_section = line + '\n'
        else:
            # Continuamos acumulando en la sección actual
            current_section += line + '\n'
    
    # Añadir la última sección si existe
    if current_section.strip():
        sections.append(current_section.strip())
    
    # Traducir cada sección
    translated_sections = []
    no_translated = 0
    
    for i, section in enumerate(sections):
        try:
            print(f"🔄 Traduciendo sección {i+1}/{len(sections)}...")
            
            # Dividir la sección en fragmentos más pequeños si es necesario
            max_chunk_size = 480
            chunks = [section[i:i + max_chunk_size] for i in range(0, len(section), max_chunk_size)]
            
            translated_section = ""
            for chunk in chunks:
                max_retries = 3
                translated_chunk = None
                
                for attempt in range(max_retries):
                    try:
                        # Codificar el texto para la URL
                        from urllib.parse import quote
                        encoded_chunk = quote(chunk)
                        
                        url = f"https://api.mymemory.translated.net/get?q={encoded_chunk}&langpair=es|en"
                        response = requests.get(url)
                        if response.status_code != 200:
                            raise Exception(f"HTTP {response.status_code}: {response.text}")
                        
                        data = response.json()
                        translated_chunk = data.get("responseData", {}).get("translatedText", "")
                        
                        if translated_chunk:
                            translated_section += translated_chunk + " "
                            print(f"✅ Fragmento traducido: {translated_chunk[:80]}...")
                            break
                        else:
                            raise Exception("Respuesta vacía o inesperada")
                    
                    except Exception as e:
                        if attempt < max_retries - 1:
                            print(f"⏳ Reintentando traducción ({attempt+1})... Error: {e}")
                            time.sleep(2)
                        else:
                            print(f"❌ Error permanente traduciendo fragmento: {e}")
                            # Conservar el texto original si hay un error permanente
                            translated_section += chunk + " "
                
                # Esperar un poco entre fragmentos para evitar limitaciones de la API
                time.sleep(0.5)
            
            # Verificar si la traducción se realizó correctamente
            translated_section = translated_section.strip()
            if translated_section.lower() == section.lower():
                print(f"⚠️ Traducción sospechosamente igual al original para la sección {i+1}")
                no_translated += 1
            
            # Preservar formato original de algunas palabras clave
            # Reemplazar traducciones comunes que podrían causar problemas
            translations = {
                "Topic:": "Tema:",
                "Question:": "Pregunta:",
                "Correct answer:": "Respuesta correcta:",
                "Relevant:": "Relevante:",
                "Quality:": "Calidad:",
                "Correctly classified:": "Clasificada correctamente:",
                "Valid correct answer:": "Respuesta correcta válida:",
                "Fragment:": "Fragmento:",
                "Complexity:": "Complejidad:",
                "Type:": "Tipo:"
            }
            
            for eng, esp in translations.items():
                translated_section = translated_section.replace(eng, esp)
            
            # Asegurarse que las opciones de respuesta mantengan el formato original (a), b), etc.)
            translated_section = re.sub(r'([A-E])\)', r'\1)', translated_section, flags=re.IGNORECASE)
            
            translated_sections.append(translated_section)
            
        except Exception as e:
            print(f"❌ Error general traduciendo la sección {i+1}: {e}")
            # En caso de error general, conservar la sección original
            translated_sections.append(section)
            no_translated += 1
    
    # Combinar todas las secciones traducidas
    final_translation = '\n\n'.join(translated_sections)
    
    print(f"\n📝 Traducción completada: {len(sections) - no_translated}/{len(sections)} secciones.")
    print(f"❗ Secciones no traducidas o parcialmente traducidas: {no_translated}")
    
    return final_translation


def translate_with_deepl(text, auth_key):
    import deepl
    
    translator = deepl.Translator(auth_key)
    
    # Si el texto es un string, simplemente tradúcelo
    if isinstance(text, str):
        if not text.strip():  # Si el texto está vacío
            return ""
        try:
            result = translator.translate_text(text, target_lang="EN-US")
            return str(result)
        except Exception as e:
            print(f"Error al traducir: {e}")
            return text  # Devuelve el texto original si hay error
    
    # Si el texto es un diccionario, traduce cada valor
    elif isinstance(text, dict):
        translated_dict = {}
        for topic, content in text.items():
            if not content.strip():  # Si el contenido está vacío
                translated_dict[topic] = ""
                continue
            try:
                result = translator.translate_text(content, target_lang="EN-US")
                translated_dict[topic] = str(result)
            except Exception as e:
                print(f"Error al traducir {topic}: {e}")
                translated_dict[topic] = content  # Mantener el contenido original si hay error
        return translated_dict
    
    # Si no es ni string ni diccionario, devuelve el texto sin modificar
    else:
        return text
    

def generate_unique_filename(base_name, extension):
    """
    Genera un nombre de archivo único añadiendo un contador al nombre si el archivo ya existe.

    :param base_name: Nombre base del archivo sin extensión.
    :param extension: Extensión del archivo, incluyendo el punto (por ejemplo, ".txt" o ".xlsx").
    :return: Un nombre de archivo único.
    """
    counter = 1
    new_name = f"{base_name}{extension}"
    while os.path.exists(new_name):
        new_name = f"{base_name}_{counter}{extension}"
        counter += 1
    return new_name

def parse_matching_definitions(raw_definitions):
    """
    Procesa las definiciones para preguntas de relación de columnas
    VERSIÓN CORREGIDA: Conserva los puntos finales cuando se separa por punto y coma.
    """
    if not raw_definitions:
        return []
    
    # Si es una lista, convertir a string para procesamiento unificado
    if isinstance(raw_definitions, list):
        if len(raw_definitions) == 1:
            text = raw_definitions[0]
        else:
            # Si ya es una lista con múltiples elementos, solo limpiar espacios
            cleaned = [d.strip() for d in raw_definitions if d.strip()]
            return cleaned
    else:
        text = str(raw_definitions)
    
    # Estrategia 1: Buscar patrones A) B) C) D) E)
    pattern_letters = re.findall(r"[A-E]\)\s*(.*?)(?=\s*[A-E]\)|$)", text, re.DOTALL)
    if pattern_letters:
        cleaned = [d.strip() for d in pattern_letters if d.strip()]
        return cleaned
    
    # Estrategia 2: Buscar patrones 1) 2) 3) 4) 5)
    pattern_numbers = re.findall(r"\d+\)\s*(.*?)(?=\s*\d+\)|$)", text, re.DOTALL)
    if pattern_numbers:
        cleaned = [d.strip() for d in pattern_numbers if d.strip()]
        return cleaned
    
    # Estrategia 3: Separar por punto y coma (ESTA ES LA LÍNEA CLAVE A CAMBIAR)
    if ';' in text:
        # CAMBIO: En lugar de re.split(r';\s*', text), usar split normal pero agregar puntos
        parts = text.split(';')
        cleaned = []
        for p in parts:
            p_clean = p.strip()
            if p_clean:
                # CAMBIO CRÍTICO: Agregar punto final si no lo tiene
                if not p_clean.endswith('.'):
                    p_clean += '.'
                cleaned.append(p_clean)
        return cleaned
    
    # Estrategia 4: Separar por comas (solo si hay múltiples comas)
    if text.count(',') >= 2:
        parts = text.split(',')
        cleaned = []
        for p in parts:
            p_clean = p.strip()
            if len(p_clean) > 10:
                # Si no termina con punto, agregarlo
                if not p_clean.endswith('.'):
                    p_clean += '.'
                cleaned.append(p_clean)
        return cleaned
    
    # Estrategia 5: Separar por puntos (solo si son frases largas)
    if text.count('.') >= 2:
        # Usar una expresión regular más inteligente
        parts = re.split(r'\.\s+(?=[A-E]\))', text)
        if len(parts) > 1:
            cleaned = []
            for i, p in enumerate(parts):
                p_clean = p.strip()
                if len(p_clean) > 20:
                    # Restaurar el punto excepto en el último elemento si ya lo tiene
                    if i < len(parts) - 1 and not p_clean.endswith('.'):
                        p_clean += '.'
                    cleaned.append(p_clean)
            return cleaned
        else:
            # Método original pero con restauración de puntos
            parts = re.split(r'\.\s+', text)
            cleaned = []
            for i, p in enumerate(parts):
                p_clean = p.strip()
                if len(p_clean) > 20:
                    if i < len(parts) - 1 or not p_clean.endswith('.'):
                        p_clean += '.'
                    cleaned.append(p_clean)
            return cleaned
    
    # Si no se puede dividir, retornar como una sola definición
    return [text.strip()]

from openpyxl.utils import get_column_letter
def escribir_respuestas_en_excel(worksheet, row, question_type, options, correct_answer, format_type="quizz"):
    """
    Escribe respuestas correctas y opciones según el formato especificado (pruebas o quizz).
    
    :param worksheet: Hoja de cálculo
    :param row: Fila actual
    :param question_type: Tipo de pregunta (e.g., Matching, Multiple Choice)
    :param options: Lista de opciones/conceptos extraídos.
    :param correct_answer: Cadena de respuesta(s) correcta(s) (ej: "Definición1, Definición2")
    :param format_type: Tipo de formato ('pruebas' o 'quizz'). Por defecto 'quizz'.
    """
    if question_type.lower() in ["matching", "relación de columnas"]:
        print(f"DEBUG_ESCRITURA: Procesando pregunta de tipo 'Matching/Relación de columnas'.")
        
        concept_start_col = 0
        definition_start_col = 0
        
        if format_type == "pruebas":
            concept_start_col = 4  # Columna D
            definition_start_col = 9 # Columna I
            
        elif format_type == "quizz":
            concept_start_col = 7  # Columna G
            definition_start_col = 13 # Columna M
            
        else:
            concept_start_col = 7  # Columna G (por defecto Quizz)
            definition_start_col = 13 # Columna M (por defecto Quizz)

        # Escribir conceptos
        # Siempre intentará escribir los primeros 5, si existen
        for idx, concepto in enumerate(options[:5]): 
            col_letter = get_column_letter(concept_start_col + idx)  
            worksheet[f"{col_letter}{row}"] = concepto
            print(f"🟦 {format_type.capitalize()} - Concepto en {col_letter}{row}: {concepto}")

        # Procesar definiciones con la función mejorada
        definiciones = parse_matching_definitions(correct_answer)
        
        # Escribir definiciones (máximo 5)
        for idx, definicion in enumerate(definiciones[:5]):
            col_letter = get_column_letter(definition_start_col + idx)
            worksheet[f"{col_letter}{row}"] = definicion
            print(f"🟩 {format_type.capitalize()} - Definición en {col_letter}{row}: {definicion}")
            
        # Validar que el número de conceptos coincida con el número de definiciones
        if len(options) != len(definiciones):
            print(f"⚠️ ADVERTENCIA: Número de conceptos ({len(options)}) no coincide con definiciones ({len(definiciones)})")
            
    else:
        print(f"DEBUG_ESCRITURA: ADVERTENCIA: escribir_respuestas_en_excel llamada para tipo '{question_type}' que NO es 'Matching/Relación de columnas'.")       


def export_txt_to_excel(txt_file, template_path, output_excel, chapter_order):
    """
    :param txt_file: Ruta al archivo .txt que contiene las preguntas generadas y validadas.
    :param template_path: Ruta al archivo de plantilla de Excel.
    :param output_excel: Ruta para guardar el archivo Excel generado.
    :param chapter_order: Lista con el orden de los capítulos.
    """
    try:
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter
        import re
        
        # Cargar la plantilla de Excel
        workbook = load_workbook(template_path)
        worksheet = workbook.active

        # Obtener el alto de fila y formato de la segunda fila (o fila base)
        default_row_height = worksheet.row_dimensions[2].height
        source_row_format = 2  # Número de la fila base para el formato

        # Leer el archivo de texto
        with open(txt_file, "r", encoding="utf-8") as file:
            content = file.read()

        print(f"DEBUG: Contenido del archivo TXT (primeros 1000 chars):")
        print(content[:1000])
        print("\n" + "="*50)

        current_topic = None
        row = 2  # Inicia en la fila 2 para no sobrescribir encabezados
        
        # Dividir por temas
        tema_sections = content.split("Tema:")
        
        for section in tema_sections:
            if not section.strip():
                continue
                
            lines = section.strip().split('\n')
            
            # Extraer el nombre del tema (primera línea después de "Tema:")
            if lines:
                current_topic = lines[0].strip()
                print(f"DEBUG: Procesando tema: {current_topic}")
            
            # Buscar preguntas en esta sección
            section_text = '\n'.join(lines)
            
            # Busca: número seguido de punto, luego "Pregunta:" y el contenido hasta la siguiente pregunta numerada o final
            question_pattern = r"(\d+)\.\s*Pregunta:\s*(.*?)(?=\d+\.\s*Pregunta:|\Z)"
            questions = re.findall(question_pattern, section_text, re.DOTALL)
            
            print(f"DEBUG: Encontradas {len(questions)} preguntas en el tema: {current_topic}")
            
            for q_num, question_content in questions:
                print(f"DEBUG: Procesando pregunta {q_num}")
                print(f"DEBUG: Contenido de la pregunta (primeros 300 chars): {question_content[:300]}")
                
                # Inicializar variables
                question_text = ""
                options = []
                question_type = ""
                correct_answer = ""
                explanation = ""
                relevante = ""
                calidad = ""
                clasificada_correctamente = ""
                respuesta_correcta_valida = ""
                complejidad = ""
                fragmento = ""
                
                
                # Procesar el contenido de la pregunta línea por línea
                question_lines = question_content.strip().split('\n')
                
                i = 0
                while i < len(question_lines):
                    line = question_lines[i].lstrip()
                    
                    # Extraer el enunciado de la pregunta (primera línea no vacía)
                    if not question_text and line and not line.startswith(("A)", "B)", "C)", "D)", "E)", "Tipo:", "Respuesta correcta:", "Relevante:", "Calidad:", "Clasificada correctamente:", "Respuesta correcta válida:", "Complejidad:", "Fragmento:")):
                        question_text = line
                        print(f"DEBUG: Enunciado encontrado: {question_text}")
                    
                    elif line.startswith(("A)", "B)", "C)", "D)", "E)")):
                        # Extraer la opción sin la letra
                        option_text = line[2:].strip()
                        options.append(option_text)
                        print(f"DEBUG: Opción encontrada: {option_text}")
                    
                    elif line.startswith("Tipo:"):
                        question_type = line.replace("Tipo:", "").strip()
                        print(f"DEBUG: Tipo encontrado: {question_type}")
                    
                    elif line.startswith("Respuesta correcta:"):
                        correct_answer = line.replace("Respuesta correcta:", "").strip()
                        print(f"DEBUG: Respuesta correcta: {correct_answer}")
                    
                    elif line.startswith("Relevante:"):
                        relevante = line.replace("Relevante:", "").strip()
                        print(f"DEBUG: Relevante: {relevante}")
                    
                    elif line.startswith("Calidad:"):
                        calidad = line.replace("Calidad:", "").strip()
                        print(f"DEBUG: Calidad: {calidad}")
                    
                    elif line.startswith("Clasificada correctamente:"):
                        clasificada_correctamente = line.replace("Clasificada correctamente:", "").strip()
                        print(f"DEBUG: Clasificada correctamente: {clasificada_correctamente}")
                    
                    elif line.startswith("Respuesta correcta válida:"):
                        respuesta_correcta_valida = line.replace("Respuesta correcta válida:", "").strip()
                        print(f"DEBUG: Respuesta correcta válida: {respuesta_correcta_valida}")
                    
                    elif line.startswith("Complejidad:"):
                        complejidad = line.replace("Complejidad:", "").strip()
                        print(f"DEBUG: Complejidad: {complejidad}")
                    
                    elif line.startswith("Fragmento:"):
                        fragmento = line.replace("Fragmento:", "").strip()
                        print(f"DEBUG: Fragmento: {fragmento}")
                    
                    i += 1
                
                # Solo procesar si tenemos al menos la pregunta
                if question_text:
                    print(f"DEBUG: Escribiendo pregunta en fila {row}")
                    
                    # Escribir en el Excel
                    worksheet[f"A{row}"] = row-1
                    worksheet[f"B{row}"] = question_type or "Multiple Choice"  # Tipo
                    worksheet[f"C{row}"] = question_text   # Pregunta       
                    
                    # Escribir opciones (columnas  D, E, F, G, H)
                    for idx, option in enumerate(options[:5]):  # Máximo 5 opciones
                        col_letter = get_column_letter(4 + idx)  # D=3, E=4, etc.
                        worksheet[f"{col_letter}{row}"] = option
                        print(f"DEBUG: Opción en columna {col_letter}: {option}")
                   
                    # Escribir respuestas correctas (columnas I,J,K,L,M)
                    if question_type.lower() in ["checkboxes"]:
                        # Escribir letras (A, C, D...) como respuestas correctas en columnas I–M
                        correct_letters = [x.strip().upper() for x in correct_answer.split(',')]
                        seen_letters = set()
                        correct_col = 9  # I = 9

                        for letter in correct_letters:
                            if letter in seen_letters:
                                continue
                            if 'A' <= letter <= 'E':  # Solo letras válidas
                                worksheet.cell(row=row, column=correct_col).value = letter
                                seen_letters.add(letter)
                                correct_col += 1
                    elif question_type.lower() in ["multiple choice", "true / false"]:
                        correct_letter = correct_answer.strip().upper()                        
                        if correct_letter in ['A', 'B', 'C', 'D', 'E']:
                            worksheet.cell(row=row, column=9).value = correct_letter  #options[pos]  # Columna I
                    elif question_type.lower() in ["matching", "relación de columnas"]:
                        escribir_respuestas_en_excel(worksheet, row, question_type, options, correct_answer, "pruebas")

                    worksheet[f"N{row}"]= current_topic  # Tema
                    # Usar los valores del TXT o valores por defecto
                    worksheet[f"O{row}"] = relevante or "Sí"  # Relevante
                    worksheet[f"P{row}"] = calidad or "5"  # Calidad
                    worksheet[f"Q{row}"] = clasificada_correctamente or "Sí"  # Clasificada correctamente
                    worksheet[f"R{row}"] = respuesta_correcta_valida or "Sí"  # Respuesta correcta válida
                    worksheet[f"S{row}"] = fragmento or "NA"  # Fragmento como fuente
                    worksheet[f"T{row}"] = complejidad or "Intermedio"  # Complejidad
                    
                    # Ajustar formato
                    worksheet.row_dimensions[row].height = default_row_height
                    
                    # Aplicar formato si existe la función
                    try:
                        apply_format(worksheet, source_row_format, row)
                    except:
                        pass  # Si no existe la función, continuar
                    
                    row += 1
                    print(f"DEBUG: Pregunta procesada, siguiente fila: {row}")
                else:
                    print(f"DEBUG: Pregunta {q_num} omitida - no se encontró enunciado")

        print(f"DEBUG: Total de preguntas procesadas: {row - 2}")

        # Guardar el archivo Excel
        workbook.save(output_excel)
        print(f"✅ Archivo Excel generado correctamente: {output_excel}")
        return True

    except Exception as e:
        import traceback
        error_msg = f"❌ Error al exportar a Excel: {e}\n{traceback.format_exc()}"
        print(error_msg)
        return False

def generate_formatted_txt_content(segments, preguntas_por_capitulo, curso, formato):
    """
    Genera el contenido del TXT en un formato limpio, numerado y compatible con la exportación a Excel.
    `preguntas_por_capitulo`: debe ser una lista de dicts con 'multiple_choice', 'verdadero_falso', 'checkboxes' por cada capítulo.
    """
    formatted_content = ""

    for idx, (topic, fragment) in enumerate(segments.items()):
        formatted_content += f"Tema: {topic}\n\n"
        counts = preguntas_por_capitulo[idx]
        mc_count = 1
        tf_count = 1
        cb_count = 1
        match_count=1

        # Multiple Choice
        for _ in range(counts.get("multiple_choice", 0)):
            formatted_content += f"{mc_count}. Pregunta: ¿Cuál de las siguientes afirmaciones sobre '{topic}' es correcta?\n"
            formatted_content += "A) Opción A\n"
            formatted_content += "B) Opción B\n"
            formatted_content += "C) Opción C\n"
            formatted_content += "D) Opción D\n"
            formatted_content += "E) Opción E\n"
            formatted_content += "Tipo: Multiple Choice\n"
            formatted_content += "Respuesta correcta: C\n"
            formatted_content += "Relevante: Sí\n"
            formatted_content += "Calidad: 5\n"
            formatted_content += "Clasificada correctamente: Sí\n"
            formatted_content += "Respuesta correcta válida: Sí\n"
            formatted_content += "Complejidad: Intermedio\n"
            formatted_content += f"Fragmento: {fragment[:200]}...\n\n"
            mc_count += 1

        # True / False
        for _ in range(counts.get("verdadero_falso", 0)):
            formatted_content += f"{tf_count}. Pregunta: '{topic}' requiere planificación estratégica.\n"
            formatted_content += "A) Verdadero\n"
            formatted_content += "B) Falso\n"
            formatted_content += "Tipo: True / False\n"
            formatted_content += "Respuesta correcta: A\n"
            formatted_content += "Relevante: Sí\n"
            formatted_content += "Calidad: 5\n"
            formatted_content += "Clasificada correctamente: Sí\n"
            formatted_content += "Respuesta correcta válida: Sí\n"
            formatted_content += "Complejidad: Básico\n"
            formatted_content += f"Fragmento: {fragment[:200]}...\n\n"
            tf_count += 1

        # Checkboxes
        for _ in range(counts.get("checkboxes", 0)):
            formatted_content += f"{cb_count}. Pregunta: ¿Cuáles de las siguientes son ventajas clave de '{topic}'?\n"
            formatted_content += "A) Alta disponibilidad\n"
            formatted_content += "B) Reducción de costos\n"
            formatted_content += "C) Escalabilidad\n"
            formatted_content += "D) Mejora en seguridad\n"
            formatted_content += "E) Eliminación de hardware\n"
            formatted_content += "Tipo: Checkboxes\n"
            formatted_content += "Respuesta correcta: A, B, C, D\n"
            formatted_content += "Relevante: Sí\n"
            formatted_content += "Calidad: 5\n"
            formatted_content += "Clasificada correctamente: Sí\n"
            formatted_content += "Respuesta correcta válida: Sí\n"
            formatted_content += "Complejidad: Intermedio\n"
            formatted_content += f"Fragmento: {fragment[:200]}...\n\n"
            cb_count += 1

         # Mathcing
        for _ in range(counts.get("matching", 0)):
            formatted_content += f"{match_count}. Pregunta: Relaciona cada práctica de {topic} con su descripción correspondiente.\n"
            formatted_content += "A) Gestión de incidentes\n"
            formatted_content += "B) Gestión de cambos\n"
            formatted_content += "C) Gestión del conocimiento\n"
            formatted_content += "D) Mesa de servicio\n"
            formatted_content += "E) Mejora continua\n"
            formatted_content += "Tipo: Matching\n"
            formatted_content += "Respuesta correcta: Restauras operación normal; Minimizar riesgos en cambios; Facilitar decisiones informadas; Punto único de contacto; Optimizar servicios\n"
            formatted_content += "Relevante: Sí\n"
            formatted_content += "Calidad: 5\n"
            formatted_content += "Clasificada correctamente: Sí\n"
            formatted_content += "Respuesta correcta válida: Sí\n"
            formatted_content += "Complejidad: Intermedio\n"
            formatted_content += f"Fragmento: {fragment[:200]}...\n\n"
            cb_count += 1
    return formatted_content

def export_txt_to_excel_quizz(txt_file, template_path, output_excel, chapter_order, format_type="quizz"):
    """
    Exporta datos a la plantilla de quizzes (.xlsm) respetando la estructura de columnas definida.
    VERSIÓN CORREGIDA que maneja correctamente las respuestas, incluyendo 'Relación de columnas'
    y soporta múltiples formatos de plantilla (pruebas/quizz).
    
    :param txt_file: Ruta al archivo de texto con las preguntas.
    :param template_path: Ruta a la plantilla de Excel (.xlsm).
    :param output_excel: Ruta donde se guardará el archivo Excel de salida.
    :param chapter_order: (Este parámetro no se usa directamente en esta función, pero se mantiene por compatibilidad).
    :param format_type: Tipo de formato de la plantilla ('pruebas' o 'quizz'). Por defecto 'quizz'.
    """
    try:
        # Cargar el archivo Excel manteniendo macros
        workbook = load_workbook(template_path, keep_vba=True)
        print("✅ Plantilla cargada:", workbook.sheetnames)
        
        if "Quiz" not in workbook.sheetnames:
            error_msg = "❌ La plantilla no contiene una hoja llamada 'Quiz'"
            log_message(error_msg)
            print(error_msg)
            return
            
        worksheet = workbook["Quiz"]
        
        try:
            if worksheet.protection.sheet:
                worksheet.protection.sheet = False
                print("🔓 Hoja 'Quiz' desprotegida")
        except Exception as e:
            print(f"⚠️ Advertencia al manipular protección: {e}")

        print("📥 Leyendo archivo de texto...")
        try:
            with open(txt_file, "r", encoding="utf-8") as file:
                content = file.read()
            print(f"✅ Archivo leído. Tamaño: {len(content)} caracteres")
        except Exception as e:
            error_msg = f"❌ Error al leer el archivo de texto: {e}"
            log_message(error_msg)
            print(error_msg)
            return

        if not content.strip():
            log_message("⚠️ El archivo de preguntas está vacío.")
            return

        current_topic = None
        topics_dict = {}
        row = 24  # Fila inicial para escribir datos
        questions_processed = 0

        lines = content.split('\n')
        i = 0
        
        while i < len(lines):
            line = lines[i].strip()
            i += 1
            
            if not line:
                continue
            
            if (line.startswith("Tema:") or 
                line.startswith("TEMA:") or 
                line.startswith("## ") or
                line.startswith("# ")):
                
                current_topic = line
                for prefix in ["Tema:", "TEMA:", "## ", "# "]:
                    if current_topic.startswith(prefix):
                        current_topic = current_topic[len(prefix):].strip()
                        break
                
                print(f"📚 Nuevo tema detectado: {current_topic}")
                if current_topic not in topics_dict:
                    topics_dict[current_topic] = 0
                continue
            
            question_patterns = [
                r"(\d+)\.\s*Pregunta:\s*(.+)",
                r"PREGUNTA:\s*(.+)",
                r"(\d+)\.\s*(.+)",
            ]
            
            question_text = None
            question_number = None
            
            for pattern in question_patterns:
                match = re.search(pattern, line, re.IGNORECASE)
                if match:
                    if len(match.groups()) == 2 and match.group(1).isdigit():
                        question_number = int(match.group(1))
                        question_text = match.group(2).strip()
                    else:
                        question_text = match.group(1).strip()
                    break
            
            if not question_text:
                continue
            
            print(f"📝 Pregunta detectada: {question_text[:50]}...")
            
            if current_topic:
                topics_dict[current_topic] += 1
                question_count = topics_dict[current_topic]
            else:
                question_count = questions_processed + 1
            
            options_raw = [] 
            correct_answer_mcq = None 
            question_type = "Multiple Choice" 
            matching_definitions_raw = [] 
            fragmento = "" 
            
            while i < len(lines):
                if i >= len(lines):
                    break
                    
                option_line = lines[i].strip()
                
                if (any(re.search(pattern, option_line, re.IGNORECASE) for pattern in question_patterns) or
                    option_line.startswith(("Tema:", "TEMA:", "## ", "# "))):
                    break
                
                i += 1
                
                if not option_line:
                    continue
                
                type_match = re.search(r"tipo\s*[:：]\s*(.+)", option_line, re.IGNORECASE)
                if type_match:
                    question_type = type_match.group(1).strip()
                    print(f"  - Tipo detectado: {question_type}")
                    continue 
                
                option_match = re.match(r"^\s*([a-eA-E])\)?\s+(.+)", option_line)
                if option_match:
                    letter = option_match.group(1).upper()
                    content = option_match.group(2).strip()
                    options_raw.append((letter, content)) 
                    print(f"  - Opción {letter}: {content[:30]}...")
                    continue
                if re.search(r"respuesta\s+correcta\s*[:：]\s*", option_line, re.IGNORECASE):
                    correct_part_raw = re.split(r"respuesta\s+correcta\s*[:：]\s*", option_line, flags=re.IGNORECASE)[1].strip()
                    
                    # Para MCQ (Multiple Choice)
                    mcq_correct_match = re.match(r"^([a-eA-E])", correct_part_raw)
                    if mcq_correct_match:
                        correct_answer_mcq = mcq_correct_match.group(1).upper()
                        print(f"  - Respuesta correcta (MCQ): {correct_answer_mcq}")
                    else:
                        # Para preguntas de relación de columnas
                        # NO procesar aquí, dejar el texto crudo para procesamiento posterior
                        matching_definitions_raw = correct_part_raw
                        print(f"  - Respuesta correcta (Matching) RAW: {matching_definitions_raw}")
                        
                        # Leer líneas adicionales si es necesario
                        while i < len(lines):
                            next_line = lines[i].strip()
                            if (re.match(r"^\s*([a-eA-E])\)?\s+(.+)", next_line) or 
                                any(re.search(pattern, next_line, re.IGNORECASE) for pattern in question_patterns) or 
                                next_line.startswith(("Tema:", "TEMA:", "## ", "# ")) or 
                                re.search(r"fragmento\s*[:：]\s*", next_line, re.IGNORECASE) or 
                                next_line.startswith(("Relevante:", "Calidad:", "Clasificada correctamente:", "Respuesta correcta válida:", "Complejidad:")) or 
                                not next_line):
                                break
                            
                            matching_definitions_raw += " " + next_line
                            i += 1
                        
                        print(f"  - Respuesta correcta (Matching) COMPLETA: {matching_definitions_raw}")
                    continue

                fragment_match = re.search(r"fragmento\s*[:：]\s*(.+)", option_line, re.IGNORECASE)
                if fragment_match:
                    fragmento = fragment_match.group(1).strip()
                    print(f"  - Fragmento detectado: {fragmento[:50]}...")
                    continue
                
                if question_type == "Multiple Choice" and re.search(r"verdadero|falso|true|false", option_line, re.IGNORECASE) and not options_raw:
                    question_type = "True / False"
                    options_raw = [("A", "Verdadero"), ("B", "Falso")]
                    print(f"  - Tipo detectado automáticamente: {question_type}")
                    continue
            
            options_for_excel = [] 
            correct_answer_for_excel = None 

            if question_type.lower() in ["matching", "relación de columnas"]:
                options_for_excel = [content for letter, content in options_raw]
                #correct_answer_for_excel = ", ".join(matching_definitions_raw)
                # Procesar las definiciones, ya sea que vengan como una sola cadena o como múltiples líneas
                if isinstance(matching_definitions_raw, list) and len(matching_definitions_raw) == 1:
                    # Puede que venga todo en una sola línea separada por A), B), etc.
                    defs_line = matching_definitions_raw[0]
                    parsed_defs = re.findall(r"[A-E]\)\s*(.*?)(?=\s*[A-E]\)|$)", defs_line)
                    if parsed_defs:
                        matching_definitions_raw = [d.strip(" ;,") for d in parsed_defs]

                correct_answer_for_excel = matching_definitions_raw



                print(f"  - Preparado para Matching: Conceptos: {options_for_excel}, Definiciones: '{correct_answer_for_excel}'")
            elif question_type == "True / False":
                if not options_raw:
                    options_for_excel = [("A", "Verdadero"), ("B", "Falso")]
                else:
                    options_for_excel = [content for letter, content in options_raw] 
                correct_answer_for_excel = correct_answer_mcq 
            else: 
                options_for_excel = [content for letter, content in options_raw] 
                correct_answer_for_excel = correct_answer_mcq 
            
            try:
                print(f"📊 Escribiendo pregunta {question_count} en fila {row}")
                
                worksheet.cell(row=row, column=1).value = question_type
                worksheet.cell(row=row, column=2).value = current_topic or "General"
                worksheet.cell(row=row, column=5).value = question_count
                worksheet.cell(row=row, column=6).value = question_text
                
                if question_type.lower() in ["matching", "relación de columnas"]:
                    print(f"DEBUG_EXPORT: Llamando a escribir_respuestas_en_excel para Matching. Tipo: '{question_type}', Len opciones: {len(options_for_excel)}")
                    
                    escribir_respuestas_en_excel(worksheet, row, question_type, options_for_excel, correct_answer_for_excel, format_type)
                    print(f"  ✅ Pregunta de Relación de Columnas colocada usando escribir_respuestas_en_excel para formato '{format_type}'.")
                    
                elif options_for_excel and correct_answer_for_excel is not None: 
                    print(f"🔍 Procesando opciones para MCQ/T/F. Respuesta correcta: {correct_answer_for_excel}")
                    
                    # Determinar la columna correcta basada en la letra de la respuesta (A->G, B->H, etc.)
                    correct_col_mcq = 7 # Por defecto a G (columna 7)
                    if correct_answer_for_excel and correct_answer_for_excel.isalpha():
                        letter_offset = ord(correct_answer_for_excel.upper()) - ord('A')
                        calculated_col = 7 + letter_offset # G es columna 7, H es 8, etc.
                        # Asegurarse de que la columna calculada esté dentro del rango G (7) a L (12)
                        if 7 <= calculated_col <= 12:
                            correct_col_mcq = calculated_col
                        else:
                            print(f"⚠️ Advertencia: La respuesta '{correct_answer_for_excel}' resulta en una columna fuera del rango G-L (7-12). Usando columna G.")
                    
                    incorrect_start_col_mcq = 13 # Columna M

                    for letter, content in options_raw: 
                        if letter == correct_answer_for_excel:
                            worksheet.cell(row=row, column=correct_col_mcq).value = content
                            print(f"  ✅ Respuesta correcta '{correct_answer_for_excel}' colocada en columna {get_column_letter(correct_col_mcq)}: {content}")
                            break
                    
                    # Colocar las respuestas incorrectas en las columnas M-P (columnas 13-16)
                    incorrect_col_current = incorrect_start_col_mcq 
                    for letter, content in options_raw: 
                        if letter != correct_answer_for_excel:
                            if incorrect_col_current <= 16:  # Máximo hasta columna P
                                worksheet.cell(row=row, column=incorrect_col_current).value = content
                                print(f"  ❌ Respuesta incorrecta '{letter}' colocada en columna {get_column_letter(incorrect_col_current)}: {content}")
                                incorrect_col_current += 1
                
                elif options_for_excel: 
                    print("⚠️ No se encontró respuesta correcta definida, asumiendo A")
                    correct_answer_for_excel = "A" # Valor por defecto si no se estableció antes
                    
                    # Aplica la misma lógica dinámica para el caso de fallback (asumiendo A)
                    correct_col_mcq = 7 
                    if correct_answer_for_excel and correct_answer_for_excel.isalpha():
                        letter_offset = ord(correct_answer_for_excel.upper()) - ord('A')
                        calculated_col = 7 + letter_offset
                        if 7 <= calculated_col <= 12:
                            correct_col_mcq = calculated_col
                        else:
                            print(f"⚠️ Advertencia: La respuesta asumida 'A' resulta en una columna fuera del rango G-L (7-12). Usando columna G.")

                    incorrect_start_col_mcq = 13 

                    for letter, content in options_raw: 
                        if letter == correct_answer_for_excel:
                            worksheet.cell(row=row, column=correct_col_mcq).value = content
                            print(f"  ✅ Respuesta correcta asumida '{correct_answer_for_excel}' colocada en columna {get_column_letter(correct_col_mcq)}: {content}")
                            break
                    
                    incorrect_col_current = incorrect_start_col_mcq
                    for letter, content in options_raw:
                        if letter != correct_answer_for_excel:
                            if incorrect_col_current <= 16:
                                worksheet.cell(row=row, column=incorrect_col_current).value = content
                                print(f"  ❌ Respuesta incorrecta '{letter}' colocada en columna {get_column_letter(incorrect_col_current)}: {content}")
                                incorrect_col_current += 1
                
                if fragmento:
                    worksheet.cell(row=row, column=19).value = fragmento 
                    print(f"📄 Fragmento colocado en columna S: {fragmento[:50]}...")
                
                questions_processed += 1
                row += 1
                print(f"✅ Pregunta {questions_processed} procesada correctamente")
                
            except Exception as e:
                print(f"❌ Error al escribir pregunta en Excel: {e}")
                import traceback
                print(traceback.format_exc())

        try:
            workbook.save(output_excel)
            print(f"✅ Archivo Excel guardado: {output_excel}")
            print(f"📊 Total de preguntas procesadas: {questions_processed}")
            
            print("📈 Resumen por tema:")
            for tema, count in topics_dict.items():
                print(f"  - {tema}: {count} preguntas")
                
        except Exception as e:
            error_msg = f"❌ Error al guardar el archivo Excel: {e}"
            print(error_msg)
            log_message(error_msg)
            
    except Exception as e:
        error_msg = f"❌ Error general en export_txt_to_excel_quizz: {e}"
        print(error_msg)
        log_message(error_msg)
        import traceback
        print(traceback.format_exc())

def calcular_preguntas_por_tipo(num_chapters, formato="Prueba"):
    if formato.lower() == 'quizz':
        min_por_capitulo=3
        max_por_capitulo=5
        
        #Generar preguntas por cada capítulo aleatoriamente entre 3 y 5
        preguntas_por_capitulo=[random.randint(min_por_capitulo,max_por_capitulo) for _ in range(num_chapters)]         
        dist_por_tipo={
            'multiple_choice':[],
            'verdadero_falso':[],
            'relacion_col':[]    
        }
        for total_preg in preguntas_por_capitulo:
            om=max(1,int(total_preg*.6))
            vf=max(1,int(total_preg*.35))
            comp=total_preg-om-vf
            dist_por_tipo['multiple_choice'].append(om)
            dist_por_tipo['verdadero_falso'].append(vf)
            dist_por_tipo['relacion_col'].append(comp)

        return{
            'detalle_por_capitulo':preguntas_por_capitulo,
            'multiple_choice':dist_por_tipo['multiple_choice'],
            'verdadero_falso':dist_por_tipo['verdadero_falso'],
            'relacion_col':dist_por_tipo['relacion_col']
        }               

    else: #formato prueba 
        min_total = 40
        max_total = 50
        min_por_capitulo = 3
        max_por_capitulo = 5

        if num_chapters == 0:
            return {
                'total': 0,
                'multiple_choice': 0,
                'verdadero_falso': 0,
                'checkboxes': 0,
                'matching':0
            }

        # Generar un total aleatorio dentro del rango global
        total_preguntas = random.randint(min_total, max_total)

        # Si hay muy pocos capítulos, romper la regla de 5 por capítulo
        # Lo resolveremos en la distribución, no aquí

        # Distribución por tipo (60% opción múltiple, 25% verdadero/falso, 15% checkboxes)
        preguntas_multiple_choice = max(1, int(total_preguntas * 0.6))
        preguntas_verdadero_falso = max(1, int(total_preguntas * 0.25))
        preguntas_matching= max(1,int(total_preguntas*.10))
        preguntas_check = total_preguntas - preguntas_multiple_choice - preguntas_verdadero_falso-preguntas_matching

        # Corrección si hay números negativos
        if preguntas_check < 1:
            preguntas_check = 1
            restante = total_preguntas - preguntas_check
            preguntas_multiple_choice = max(1, int(restante * 0.7))
            preguntas_verdadero_falso = restante - preguntas_multiple_choice

        return {
            'total': total_preguntas,
            'multiple_choice': preguntas_multiple_choice,
            'verdadero_falso': preguntas_verdadero_falso,
            'checkboxes': preguntas_check,
            'matching': preguntas_matching
        }

        
def get_file_content(uploaded_files=None, urls=None):
    """
    Extrae contenido de archivos subidos y/o URLs remotas.

    :param uploaded_files: Lista de archivos subidos (Streamlit UploadedFile).
    :param urls: Lista de enlaces remotos a descargar.
    :return: Contenido combinado de todos los archivos.
    """
    combined_content = ""

    # Procesar archivos subidos localmente
    if uploaded_files:
        for uploaded_file in uploaded_files:
            file_name = uploaded_file.name  # Obtiene el nombre del archivo
            file_extension = os.path.splitext(uploaded_file.name)[1].lower()
            try:
                if file_extension == ".pdf":
                    with pdfplumber.open(BytesIO(uploaded_file.read())) as pdf:
                        for page in pdf.pages:
                            combined_content += f"source: {file_name}\n" + page.extract_text() + "\n"
                    log_message(f"✅ PDF cargado correctamente: {file_name}")
                elif file_extension in [".xls", ".xlsx"]:
                    workbook = openpyxl.load_workbook(BytesIO(uploaded_file.read()))
                    for sheet in workbook.sheetnames:
                        sheet_obj = workbook[sheet]
                        for row in sheet_obj.iter_rows(values_only=True):
                            combined_content += f"source: {file_name}\n" + "\t".join([str(cell) for cell in row if cell]) + "\n"
                    log_message(f"✅ Excel cargado correctamente: {file_name}")
                elif file_extension == ".pptx":
                    presentation = Presentation(BytesIO(uploaded_file.read()))
                    for idx, slide in enumerate(presentation.slides, start=1):
                        combined_content += f"source: {file_name}\nDiapositiva {idx}:\n"
                        # Extraer texto de las formas (texto en la diapositiva)
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                combined_content += shape.text + "\n"
                        # Extraer texto de las notas (si existen)
                        if slide.notes_slide and slide.notes_slide.notes_text_frame:
                            combined_content += "\nNotas de la diapositiva:\n"
                            combined_content += slide.notes_slide.notes_text_frame.text + "\n"
                        combined_content += "\n" + "-" * 50 + "\n"  # Separador entre diapositivas
                    log_message(f"✅ PowerPoint cargado correctamente: {file_name}")
                elif file_extension == ".docx":
                    document = Document(BytesIO(uploaded_file.read()))
                    for paragraph in document.paragraphs:
                        if paragraph.text.strip():
                            combined_content += f"source: {file_name}\n" + paragraph.text + "\n"
                    log_message(f"✅ Word cargado correctamente: {file_name}")
                elif file_extension == ".txt":
                    combined_content += f"source: {file_name}\n" + uploaded_file.read().decode("utf-8") + "\n"
                    log_message(f"✅ TXT cargado correctamente: {file_name}")

                else:
                    error_msg = f"❌ Formato de archivo no soportado: {file_name}"
                    log_message(error_msg)
                    st.error(error_msg)
                
            except Exception as e:
                error_msg = f"❌ Error al cargar el archivo {file_name}: {e}"
                log_message(error_msg)
                st.error(error_msg)

    # Procesar archivos desde URLs
    if urls:
        for url in urls:
            if not url.strip():
                continue
            try:
                response = requests.get(url.strip(), headers={"User-Agent": "Mozilla/5.0"})
                response.raise_for_status()
                file_extension = os.path.splitext(url.strip())[1].lower()
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=file_extension)
                temp_file.write(response.content)
                temp_file.close()
                
                # Procesar según el tipo de archivo
                if file_extension == ".pdf":
                    with pdfplumber.open(temp_file.name) as pdf:
                        for page in pdf.pages:
                            combined_content += f"source: {url}\n" + page.extract_text() + "\n"
                    log_message(f"✅ PDF descargado correctamente desde URL: {url}")
                elif file_extension in [".xls", ".xlsx"]:
                    workbook = openpyxl.load_workbook(temp_file.name)
                    for sheet in workbook.sheetnames:
                        sheet_obj = workbook[sheet]
                        for row in sheet_obj.iter_rows(values_only=True):
                            combined_content += f"source: {url}\n" + "\t".join([str(cell) for cell in row if cell]) + "\n"
                    log_message(f"✅ EXCEL descargado correctamente desde URL: {url}")
                elif file_extension == ".pptx":
                    presentation = Presentation(temp_file.name)
                    for idx, slide in enumerate(presentation.slides, start=1):
                        combined_content += f"source: {url}\nDiapositiva {idx}:\n"
                        # Extraer texto de las formas (texto en la diapositiva)
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                combined_content += shape.text + "\n"
                        # Extraer texto de las notas (si existen)
                        if slide.notes_slide and slide.notes_slide.notes_text_frame:
                            combined_content += "\nNotas de la diapositiva:\n"
                            combined_content += slide.notes_slide.notes_text_frame.text + "\n"
                        combined_content += "\n" + "-" * 50 + "\n"  # Separador entre diapositivas
                    log_message(f"✅ POWER POINT descargado correctamente desde URL: {url}")
                elif file_extension == ".docx":
                    document = Document(temp_file.name)
                    for paragraph in document.paragraphs:
                        if paragraph.text.strip():
                            combined_content += f"source: {url}\n" + paragraph.text + "\n"
                    log_message(f"✅ WORD descargado correctamente desde URL: {url}")
                elif file_extension == ".txt":
                    with open(temp_file.name, "r", encoding="utf-8") as file:
                        combined_content += f"source: {url}\n" + file.read() + "\n"
                    log_message(f"✅ TXT descargado correctamente desde URL: {url}")
                else:
                    print(f"❌ Formato no soportado desde URL: {url}")
                    log_message(f"❌ Formato no soportado desde URL: {url}")
                    st.error(f"❌ Formato no soportado desde URL: {url}")
            
                os.unlink(temp_file.name)  # Eliminar archivo temporal

            except Exception as e:
                print(f"Error al descargar o procesar URL {url}: {e}")
                error_msg = f"❌ Error al descargar archivo desde URL {url}: {e}"
                log_message(error_msg)
                st.error(error_msg)
    
    return combined_content

def get_username():
    """Obtiene solo el nombre del usuario."""
    try:
        return getpass.getuser() or os.getenv('USERNAME') or os.getenv('USER') or 'Usuario_Desconocido'
    except Exception:
        return 'Usuario_Desconocido'

def log_message(message):
    """Registra mensajes con el nombre del usuario y la marca de tiempo."""
    log_dir=os.path.join(os.getenv("LOCALAPPDATA"),"GeneradorPreguntas")
    #log_dir = "Archivos"
    os.makedirs(log_dir, exist_ok=True)
    
    # Usar un archivo de log fijo (NO generar uno nuevo cada vez)
    log_path = os.path.join(log_dir, "log.txt")

    # Obtener el nombre del usuario
    usuario = get_username()

    # Escribir el mensaje en el log
    with open(log_path, "a", encoding="utf-8") as log_file:
        timestamp = time.strftime("%Y-%m-%d %H:%M:%S")
        log_file.write(f"Usuario: {usuario} [{timestamp}] {message}\n")

def global_exception_handler(exctype, value, tb):
    """
    Captura errores globales no manejados y los escribe en el log.
    """
    error_message = "".join(traceback.format_exception(exctype, value, tb))
    log_message(f"❌ ERROR CRÍTICO NO MANEJADO:\n{error_message}")
    st.error("❌ Se produjo un error crítico. Revisa el archivo de logs para más detalles.")

def check_dependencies():
    """
    Verifica si las librerías críticas están actualizadas.
    """
    import pkg_resources

    required_packages = ["streamlit", "openpyxl", "pdfplumber", "requests", "openai"]

    for package in required_packages:
        try:
            pkg_resources.get_distribution(package)
            log_message(f"✅ {package} está instalada.")
        except pkg_resources.DistributionNotFound:
            log_message(f"❌ La librería {package} no está instalada.")
            st.error(f"La librería {package} no está instalada.")
        except Exception as e:
            log_message(f"⚠️ Error al verificar {package}: {e}")
